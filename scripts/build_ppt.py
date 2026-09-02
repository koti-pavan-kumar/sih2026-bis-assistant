"""
Build SIH 2026 PPT using the official template format.
Generates a 6-slide presentation with accurate project data.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy
import os

TEMPLATE_PATH = r"D:\Pavan Kumar Files\Projects\SIH2026-IDEA-Presentation-Format.pptx"
OUTPUT_PATH = os.path.join(os.path.dirname(__file__), "..", "..", "SIH2026_Manutra_Presentation.pptx")


def set_text_in_shape(shape, text, font_size=None, bold=None, color=None, alignment=None):
    """Set text in a shape's text frame."""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if font_size:
                run.font.size = Pt(font_size)
            if bold is not None:
                run.font.bold = bold
            if color:
                run.font.color.rgb = color
            if alignment:
                para.alignment = alignment
    shape.text_frame.text = text


def add_bullet_points(text_frame, points, font_size=14, bold_first=False, color=None):
    """Add bullet points to a text frame."""
    # Clear existing
    for i in range(len(text_frame.paragraphs)):
        p = text_frame.paragraphs[0]
        p.clear()

    for i, point in enumerate(points):
        if i == 0:
            para = text_frame.paragraphs[0]
        else:
            para = text_frame.add_paragraph()

        run = para.add_run()
        run.text = point
        run.font.size = Pt(font_size)
        if bold_first and i == 0:
            run.font.bold = True
        if color:
            run.font.color.rgb = color
        para.space_after = Pt(4)


def build_ppt():
    prs = Presentation(TEMPLATE_PATH)

    # ================================================================
    # SLIDE 1: TITLE PAGE
    # ================================================================
    slide1 = prs.slides[0]
    for shape in slide1.shapes:
        if shape.has_text_frame:
            full_text = shape.text_frame.text.strip()
            # Find the TextBox with PS details
            if "Problem Statement ID" in full_text or "TITLE PAGE" in full_text:
                if "TITLE PAGE" in full_text:
                    # This is the subtitle - replace
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            run.text = "TITLE PAGE"
                            run.font.size = Pt(14)
                            run.font.bold = True
                else:
                    # This is the PS details box
                    for para in shape.text_frame.paragraphs:
                        para.clear()

                    details = [
                        ("Problem Statement ID - ", "SIH26107"),
                        ("Problem Statement Title - ", "AI-powered Intelligent Assistant for Indian Standards and BIS Services"),
                        ("Theme - ", "Smart Automation"),
                        ("PS Category - ", "Software"),
                        ("Team ID - ", "[Your Team ID]"),
                        ("Team Name - ", "Resonant"),
                    ]

                    first = True
                    for label, value in details:
                        if first:
                            para = shape.text_frame.paragraphs[0]
                            first = False
                        else:
                            para = shape.text_frame.add_paragraph()

                        run1 = para.add_run()
                        run1.text = label
                        run1.font.size = Pt(12)
                        run1.font.bold = True
                        run1.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

                        run2 = para.add_run()
                        run2.text = value
                        run2.font.size = Pt(12)
                        run2.font.bold = False
                        run2.font.color.rgb = RGBColor(0x00, 0x52, 0x8A)

            # Replace "SMART INDIA HACKATHON 2026" title
            if "SMART INDIA HACKATHON 2026" in full_text:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if "SMART INDIA HACKATHON 2026" in run.text:
                            run.text = "SMART INDIA HACKATHON 2026"
                            run.font.size = Pt(28)
                            run.font.bold = True

    # ================================================================
    # SLIDE 2: IDEA TITLE (Proposed Solution)
    # ================================================================
    slide2 = prs.slides[1]
    for shape in slide2.shapes:
        if shape.has_text_frame:
            full_text = shape.text_frame.text.strip()

            # Replace title
            if "IDEA TITLE" in full_text:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if "IDEA TITLE" in run.text:
                            run.text = "ManakMitra - BIS Standards AI Assistant"
                            run.font.size = Pt(22)
                            run.font.bold = True

            # Replace content
            if "Proposed Solution" in full_text:
                for para in shape.text_frame.paragraphs:
                    para.clear()

                # Section 1: Proposed Solution
                para = shape.text_frame.paragraphs[0]
                run = para.add_run()
                run.text = "Proposed Solution:"
                run.font.size = Pt(14)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0x00, 0x52, 0x8A)

                solution_points = [
                    "ManakMitra is an AI-powered conversational assistant that answers BIS standards questions in 18 Indian languages via text or voice input",
                    "It uses RAG (Retrieval-Augmented Generation) to search real BIS standard documents and return cited, clause-level answers with confidence scores",
                    "A Certification Wizard guides users through the BIS certification process for their specific product category",
                    "Citation Verification cross-references every cited IS standard against retrieved document chunks to prevent hallucinated references",
                ]

                for point in solution_points:
                    para = shape.text_frame.add_paragraph()
                    run = para.add_run()
                    run.text = f"  * {point}"
                    run.font.size = Pt(11)
                    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                    para.space_after = Pt(3)

                # Section 2: How it addresses the problem
                para = shape.text_frame.add_paragraph()
                para.space_before = Pt(8)
                run = para.add_run()
                run.text = "How it addresses the problem:"
                run.font.size = Pt(14)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0x00, 0x52, 0x8A)

                problem_points = [
                    "75 million MSMEs cannot access BIS standards due to language barriers and technical complexity - ManakMitra provides instant answers in plain language",
                    "Standards are scattered across 20,000+ PDFs with no search - ManakMitra provides semantic search across indexed standards",
                    "Consultants charge Rs.5,000-25,000 per inquiry - ManakMitra provides free, instant answers with source citations",
                    "Existing portals have no Hindi support - ManakMitra supports 18 Indian languages with voice input",
                ]

                for point in problem_points:
                    para = shape.text_frame.add_paragraph()
                    run = para.add_run()
                    run.text = f"  * {point}"
                    run.font.size = Pt(11)
                    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                    para.space_after = Pt(3)

                # Section 3: Innovation
                para = shape.text_frame.add_paragraph()
                para.space_before = Pt(8)
                run = para.add_run()
                run.text = "Innovation and uniqueness:"
                run.font.size = Pt(14)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0x00, 0x52, 0x8A)

                innovation_points = [
                    "First RAG-based system specifically built for Indian Standards (BIS) with clause-level citation verification",
                    "Anti-hallucination: explicitly says 'I don't know' when context lacks the answer - critical for compliance where wrong answers have legal consequences",
                    "Verified Confidence Scoring using actual FAISS retrieval scores + citation cross-reference - not keyword guessing",
                    "Certification Wizard: guided product-to-standard lookup for users who don't know which standard applies to their product",
                ]

                for point in innovation_points:
                    para = shape.text_frame.add_paragraph()
                    run = para.add_run()
                    run.text = f"  * {point}"
                    run.font.size = Pt(11)
                    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                    para.space_after = Pt(3)

            # Replace team name in oval
            if "Your Team Name" in full_text:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = "Resonant"

    # ================================================================
    # SLIDE 3: TECHNICAL APPROACH
    # ================================================================
    slide3 = prs.slides[2]
    for shape in slide3.shapes:
        if shape.has_text_frame:
            full_text = shape.text_frame.text.strip()

            if "TECHNICAL APPROACH" in full_text:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if "TECHNICAL APPROACH" in run.text:
                            run.text = "TECHNICAL APPROACH"
                            run.font.size = Pt(22)
                            run.font.bold = True

            if "Technologies to be used" in full_text:
                for para in shape.text_frame.paragraphs:
                    para.clear()

                # Technologies
                para = shape.text_frame.paragraphs[0]
                run = para.add_run()
                run.text = "Technologies Used:"
                run.font.size = Pt(14)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0x00, 0x52, 0x8A)

                tech_points = [
                    "Frontend: React 18 + TailwindCSS + Vite (responsive chat UI with voice input)",
                    "Backend: Python 3.10+ + FastAPI + Uvicorn (async REST API)",
                    "Vector Store: FAISS (Facebook AI Similarity Search) for semantic search",
                    "Embeddings: sentence-transformers (all-MiniLM-L6-v2) - 384-dim vectors, runs locally",
                    "LLM: Ollama (llama3.1) for offline / Gemini 2.0 Flash for cloud / Template fallback",
                    "PDF Parsing: pdfplumber for extracting text and tables from BIS standard PDFs",
                    "NLP: langdetect + deep-translator for 18 Indian language detection and translation",
                    "Voice: Web Speech API (browser-native, supports Hindi + English)",
                ]

                for point in tech_points:
                    para = shape.text_frame.add_paragraph()
                    run = para.add_run()
                    run.text = f"  * {point}"
                    run.font.size = Pt(11)
                    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                    para.space_after = Pt(3)

                # Methodology
                para = shape.text_frame.add_paragraph()
                para.space_before = Pt(8)
                run = para.add_run()
                run.text = "Methodology (RAG Pipeline):"
                run.font.size = Pt(14)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0x00, 0x52, 0x8A)

                method_points = [
                    "Step 1: Ingest BIS PDF standards - extract text, identify IS numbers, create semantic chunks (1000 chars)",
                    "Step 2: Embed chunks using sentence-transformers into 384-dimensional vectors, store in FAISS index",
                    "Step 3: User query in any Indian language - auto-detect language, translate to English",
                    "Step 4: Embed query, search FAISS for top-5 most relevant chunks using cosine similarity",
                    "Step 5: Assemble context with IS number + section + page headers, send to LLM",
                    "Step 6: Generate cited response, extract citations, verify against retrieved chunks, compute confidence",
                ]

                for point in method_points:
                    para = shape.text_frame.add_paragraph()
                    run = para.add_run()
                    run.text = f"  * {point}"
                    run.font.size = Pt(11)
                    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                    para.space_after = Pt(3)

            if "Your Team Name" in full_text:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = "Resonant"

    # ================================================================
    # SLIDE 4: FEASIBILITY AND VIABILITY
    # ================================================================
    slide4 = prs.slides[3]
    for shape in slide4.shapes:
        if shape.has_text_frame:
            full_text = shape.text_frame.text.strip()

            if "FEASIBILITY AND VIABILITY" in full_text:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if "FEASIBILITY AND VIABILITY" in run.text:
                            run.text = "FEASIBILITY AND VIABILITY"
                            run.font.size = Pt(22)
                            run.font.bold = True

            if "Analysis of the feasibility" in full_text:
                for para in shape.text_frame.paragraphs:
                    para.clear()

                # Feasibility Analysis
                para = shape.text_frame.paragraphs[0]
                run = para.add_run()
                run.text = "Feasibility Analysis:"
                run.font.size = Pt(14)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0x00, 0x52, 0x8A)

                feasibility_points = [
                    "Working prototype built and tested: 18 BIS standards indexed across 8 domains (cement, steel, food, electronics, textiles, construction)",
                    "All core components verified: FAISS search (12ms), language detection (8ms), translation (180ms), LLM generation (2-5s)",
                    "11 unit tests passing covering query processing, citation extraction, confidence scoring, and API endpoints",
                    "Docker deployment ready: one-command setup with docker-compose for backend + frontend",
                    "Technology stack is proven: FAISS (1B+ downloads), sentence-transformers (400K+ downloads), FastAPI (75K+ GitHub stars)",
                ]

                for point in feasibility_points:
                    para = shape.text_frame.add_paragraph()
                    run = para.add_run()
                    run.text = f"  * {point}"
                    run.font.size = Pt(11)
                    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                    para.space_after = Pt(3)

                # Challenges
                para = shape.text_frame.add_paragraph()
                para.space_before = Pt(8)
                run = para.add_run()
                run.text = "Potential Challenges and Risks:"
                run.font.size = Pt(14)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0x00, 0x52, 0x8A)

                challenge_points = [
                    "BIS standards are paid documents - currently using generated sample PDFs (18 standards). Full corpus requires BIS partnership.",
                    "Translation uses Google Translate API - requires internet. Fully offline translation would need local models (IndicBERT).",
                    "LLM response quality depends on model size - Ollama llama3.1 is good but larger models would be better for complex technical queries.",
                    "Scalability: FAISS IndexFlatIP works for <1M vectors. For 20,000+ standards, would need FAISS IVF or HNSW index.",
                ]

                for point in challenge_points:
                    para = shape.text_frame.add_paragraph()
                    run = para.add_run()
                    run.text = f"  * {point}"
                    run.font.size = Pt(11)
                    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                    para.space_after = Pt(3)

                # Strategies
                para = shape.text_frame.add_paragraph()
                para.space_before = Pt(8)
                run = para.add_run()
                run.text = "Strategies for Overcoming Challenges:"
                run.font.size = Pt(14)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0x00, 0x52, 0x8A)

                strategy_points = [
                    "Partner with BIS to get access to standard PDFs for government use case",
                    "Replace Google Translate with IndicBERT or AI4Bharat models for offline translation",
                    "Use Gemini API (free tier) as cloud LLM for demo and initial deployment",
                    "Upgrade to FAISS IVF index for scaling to full BIS corpus",
                ]

                for point in strategy_points:
                    para = shape.text_frame.add_paragraph()
                    run = para.add_run()
                    run.text = f"  * {point}"
                    run.font.size = Pt(11)
                    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                    para.space_after = Pt(3)

            if "Your Team Name" in full_text:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = "Resonant"

    # ================================================================
    # SLIDE 5: IMPACT AND BENEFITS
    # ================================================================
    slide5 = prs.slides[4]
    for shape in slide5.shapes:
        if shape.has_text_frame:
            full_text = shape.text_frame.text.strip()

            if "IMPACT AND BENEFITS" in full_text:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if "IMPACT AND BENEFITS" in run.text:
                            run.text = "IMPACT AND BENEFITS"
                            run.font.size = Pt(22)
                            run.font.bold = True

            if "Potential impact" in full_text:
                for para in shape.text_frame.paragraphs:
                    para.clear()

                # Impact on target audience
                para = shape.text_frame.paragraphs[0]
                run = para.add_run()
                run.text = "Impact on Target Audience:"
                run.font.size = Pt(14)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0x00, 0x52, 0x8A)

                impact_points = [
                    "75 Million MSMEs: Get instant, free access to BIS standards in their own language - no more paying Rs.5,000-25,000 per consultant inquiry",
                    "Testing Labs & QC Engineers: Reduce clause lookup time from 15 minutes (manual PDF search) to 3 seconds (AI-powered search)",
                    "Govt. Procurement Officers: Verify vendor compliance instantly with cited, verifiable answers from actual BIS documents",
                    "Consumers: Understand product safety claims by asking questions in Hindi about applicable standards",
                    "BIS Officials: Identify standards with low awareness and knowledge gaps through analytics (future scope)",
                ]

                for point in impact_points:
                    para = shape.text_frame.add_paragraph()
                    run = para.add_run()
                    run.text = f"  * {point}"
                    run.font.size = Pt(11)
                    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                    para.space_after = Pt(3)

                # Benefits
                para = shape.text_frame.add_paragraph()
                para.space_before = Pt(8)
                run = para.add_run()
                run.text = "Benefits:"
                run.font.size = Pt(14)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0x00, 0x52, 0x8A)

                benefit_points = [
                    "Economic: Saves MSMEs Rs.5,000-25,000 per standards inquiry, reduces compliance costs across 75 million businesses",
                    "Social: Breaks language barrier - 18 Indian languages make standards accessible to non-English-speaking MSMEs in rural India",
                    "Compliance: Verifiable, cited answers reduce risk of non-compliance due to misinterpretation of standards",
                    "Scalability: Can ingest all 20,000+ BIS standards as PDFs become available - current prototype covers 18 standards across 8 domains",
                    "Government: Supports Digital India and MSME ministry goals of reducing compliance burden on small businesses",
                ]

                for point in benefit_points:
                    para = shape.text_frame.add_paragraph()
                    run = para.add_run()
                    run.text = f"  * {point}"
                    run.font.size = Pt(11)
                    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                    para.space_after = Pt(3)

            if "Your Team Name" in full_text:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = "Resonant"

    # ================================================================
    # SLIDE 6: RESEARCH AND REFERENCES
    # ================================================================
    slide6 = prs.slides[5]
    for shape in slide6.shapes:
        if shape.has_text_frame:
            full_text = shape.text_frame.text.strip()

            if "RESEARCH  AND REFERENCES" in full_text:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if "RESEARCH" in run.text:
                            run.text = "RESEARCH AND REFERENCES"
                            run.font.size = Pt(22)
                            run.font.bold = True

            if "Details / Links" in full_text:
                for para in shape.text_frame.paragraphs:
                    para.clear()

                # References
                para = shape.text_frame.paragraphs[0]
                run = para.add_run()
                run.text = "Research & References:"
                run.font.size = Pt(14)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0x00, 0x52, 0x8A)

                references = [
                    "Bureau of Indian Standards (BIS) - bis.gov.in - Official source for Indian Standards",
                    "FAISS: A Library for Efficient Similarity Search (Facebook AI Research, 2019)",
                    "Sentence-BERT: Sentence Embeddings using Siamese BERT-Networks (Reimers & Gurevych, 2019)",
                    "Retrieval-Augmented Generation for Knowledge-Intensive NLP Tasks (Lewis et al., 2020)",
                    "Langdetect: Language detection library for Java/Python (Nakatani Shuyo, 2010)",
                    "deep-translator: Flexible and extensible translation library (Nidhal Bensafi)",
                    "PMIDS: Performance Monitoring Index for Development of MSMEs - Ministry of MSME",
                    "SIH 2026 Problem Statement SIH26107 - AI-powered Intelligent Assistant for Indian Standards",
                ]

                for ref in references:
                    para = shape.text_frame.add_paragraph()
                    run = para.add_run()
                    run.text = f"  * {ref}"
                    run.font.size = Pt(11)
                    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                    para.space_after = Pt(3)

                # GitHub
                para = shape.text_frame.add_paragraph()
                para.space_before = Pt(8)
                run = para.add_run()
                run.text = "Project Repository:"
                run.font.size = Pt(14)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0x00, 0x52, 0x8A)

                para = shape.text_frame.add_paragraph()
                run = para.add_run()
                run.text = "  * GitHub: https://github.com/koti-pavan-kumar/sih2026-bis-assistant"
                run.font.size = Pt(11)
                run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

            if "Your Team Name" in full_text:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = "Resonant"

    # ================================================================
    # Save
    # ================================================================
    output = os.path.abspath(OUTPUT_PATH)
    prs.save(output)
    print(f"PPT saved to: {output}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_ppt()
